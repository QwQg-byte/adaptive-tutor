# -*- coding: utf-8 -*-
"""扫描算法设计课件目录，提取每章小节结构，输出 JSON 供人工整理知识点清单"""
import argparse
import json
import re
import sys
from pathlib import Path
from pptx import Presentation

for s in (sys.stdout, sys.stderr):
    try:
        s.reconfigure(encoding='utf-8')
    except Exception:
        pass

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SOURCE = PROJECT_ROOT / 'data' / 'raw' / '算法设计课件'

parser = argparse.ArgumentParser(description=__doc__)
parser.add_argument('--source', type=Path, default=DEFAULT_SOURCE, help='课件目录')
parser.add_argument('--output', type=Path, help='输出 JSON，默认写到课件目录下')
args = parser.parse_args()
SRC = args.source.resolve()
OUT = (args.output or (SRC / '_scan_result.json')).resolve()
if not SRC.is_dir():
    parser.error(f'课件目录不存在: {SRC}')

chap_pat = re.compile(r'第(\d{1,2})章')
result = {}

for ppt_file in sorted(SRC.glob('*.pptx')):
    m = chap_pat.search(ppt_file.stem)
    if not m:
        continue
    chap_no = int(m.group(1))
    chap_name = chap_pat.sub('', ppt_file.stem).strip()

    slides_text = []
    try:
        prs = Presentation(str(ppt_file))
    except Exception as e:
        print(f'[跳过] {ppt_file.name}: {e}')
        continue
    for slide in prs.slides:
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
        slides_text.append(texts)

    # 提取 N.M 形式的小节标题（一级小节，不含 N.M.K）
    sec_pat = re.compile(rf'^{chap_no}\.(\d{{1,2}})(?!\.)\*?\s*(\S.*)$')
    sections = {}
    for texts in slides_text:
        for t in texts:
            for line in t.split('\n'):
                line = re.sub(r'\s+', ' ', line.strip())
                m2 = sec_pat.match(line)
                if m2:
                    n = int(m2.group(1))
                    title = m2.group(2).strip()
                    # 去掉尾部页码等噪声；取最短的干净标题
                    title = re.sub(r'[0-9]+$', '', title).strip()
                    if title and (n not in sections or len(title) < len(sections[n])):
                        sections[n] = title

    result[chap_no] = {
        'file': ppt_file.name,
        'chapter': chap_name,
        'slides': len(slides_text),
        'sections': {f'{chap_no}.{n}': sections[n] for n in sorted(sections)},
        # 前3页文本帮助判断章节主旨
        'intro_preview': [t[:120] for texts in slides_text[:3] for t in texts][:8],
    }

with OUT.open('w', encoding='utf-8') as output_file:
    json.dump(result, output_file, ensure_ascii=False, indent=2)
for no in sorted(result):
    r = result[no]
    print(f"第{no}章 {r['chapter']}  ({r['slides']}页, {len(r['sections'])}节)")
    for k, v in r['sections'].items():
        print(f'   {k} {v}')
print(f'\n[输出] {OUT}')
